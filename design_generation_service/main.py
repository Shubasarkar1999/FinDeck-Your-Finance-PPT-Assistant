# main.py

# --- Imports ---
import os
import uuid
import logging
import tempfile
import base64
import io
import random
from typing import Dict, List, Tuple, Optional

import httpx
from fastapi import FastAPI, HTTPException
from google.cloud import storage
from pptx import Presentation
from pptx.slide import Slide as PptxSlide
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER
import urllib.parse

# Import your models from models.py
from models import GenerationRequest, GenerationResponse, ImageServiceRequest, Slide

# --- Configuration ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - [%(levelname)s] - %(message)s')
logger = logging.getLogger(__name__)
app = FastAPI(title="Design & Generation Service (python-pptx)")
SERVICE_ACCOUNT_KEY_PATH = "sunlit-runway-472202-p8-75230f6c1db6.json"
BUCKET_NAME = "finance-ppt-bot"
IMAGE_SERVICE_URL = os.environ.get("IMAGE_SERVICE_URL")
storage_client = storage.Client.from_service_account_json(SERVICE_ACCOUNT_KEY_PATH)

THEME_MAP = {
    "minimalist": "templates/Dark.pptx",
    "streamline": "templates/Streamline.pptx",
    "forest": "templates/Forest.pptx",
    "sunset": "templates/Sunset.pptx",
    "orbit": "templates/Orbit.pptx",
    "mystique": "templates/Mystique.pptx"
}
DEFAULT_TEMPLATE = "templates/Dark.pptx"

# --- UPDATED: New layouts added ---
LAYOUT_MAP = {
    "title_slide": 0, "bullet_points": 1, "image_left": 2,
    "image_right": 3, "conclusion_slide": 0,
    # --- ADDITIONS START HERE ---
    "sticker_left": 4,
    "sticker_right": 5
    # --- ADDITIONS END HERE ---
}

# --- Helper Functions ---
def get_placeholder(slide: PptxSlide, ph_type: PP_PLACEHOLDER) -> Optional[SlidePlaceholder]:
    """
    Finds a placeholder on a slide by its type.
    Returns the placeholder shape or None if not found.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type == ph_type:
            return shape
    return None

def upload_to_gcs(source_file_path: str, destination_blob_name: str, job_id: str) -> str:
    bucket = storage_client.bucket(BUCKET_NAME)
    blob = bucket.blob(destination_blob_name)
    blob.upload_from_filename(source_file_path)
    logger.info(f"[{job_id}] Upload complete. URL: {blob.public_url}")
    return blob.public_url

def strategically_add_image_layouts(slides: List[Slide]) -> List[Slide]:
    if len(slides) <= 2: return slides
    content_slides_indices = [i for i, s in enumerate(slides) if s.layout == "bullet_points"]
    if not content_slides_indices: return slides
    num_slides_to_change = int(len(content_slides_indices) * 0.8)
    if num_slides_to_change == 0 and len(content_slides_indices) > 0: num_slides_to_change = 1
    indices_to_change = random.sample(content_slides_indices, num_slides_to_change)
    for i in indices_to_change:
        slides[i].layout = random.choice(["image_left", "image_right"])
    return slides

# --- NEW: Function to add sticker layouts for visual variety ---
def strategically_add_sticker_layouts(slides: List[Slide]) -> List[Slide]:
    """
    Takes slides that are still 'bullet_points' and randomly changes some
    to sticker layouts.
    """
    # Find slides that are still standard bullet points after the main image function has run
    bullet_point_indices = [i for i, s in enumerate(slides) if s.layout == "bullet_points"]
    if not bullet_point_indices:
        return slides

    # Decide how many sticker slides to add (e.g., 1 or 2 for variety)
    num_stickers_to_add = min(len(bullet_point_indices), 2)
    
    # Randomly select and change them
    indices_to_change = random.sample(bullet_point_indices, num_stickers_to_add)
    for i in indices_to_change:
        slides[i].layout = random.choice(["sticker_left", "sticker_right"])
        
    return slides

# --- UPDATED: Now identifies sticker layouts as needing images ---
def identify_slides_for_imaging(slides: List[Slide]) -> Tuple[List[Slide], Dict[int, int]]:
    slides_needing_images, index_map = [], {}
    layouts_needing_images = ["image_left", "image_right", "sticker_left", "sticker_right"]
    for i, slide in enumerate(slides):
        if slide.layout in layouts_needing_images:
            index_map[i] = len(slides_needing_images)
            slides_needing_images.append(slide)
    return slides_needing_images, index_map

def generate_html_preview(prs: Presentation) -> str:
    html = ['<html><head><style>body{font-family:sans-serif;padding:2rem;} .slide{border:1px solid #ccc; padding:1rem; margin-bottom:1rem; border-radius:5px;}</style></head><body>']
    html.append("<h1>Presentation Content Preview</h1>")
    for i, slide in enumerate(prs.slides):
        html.append(f'<div class="slide"><h2>Slide {i+1}</h2>')
        for shape in slide.shapes:
            if shape.has_text_frame:
                html.append(f"<p>{shape.text.replace('**', '')}</p>")
        html.append('</div>')
    html.append('</body></html>')
    return "".join(html)

# --- Main Endpoint ---
@app.post("/generate-full-presentation", response_model=GenerationResponse)
async def generate_full_presentation(request: GenerationRequest):
    job_id = str(uuid.uuid4())
    logger.info(f"[{job_id}] Received new presentation request with theme: '{request.theme}'.")

    request.slides = [s for s in request.slides if s.data and ((s.data.title and s.data.title.strip()) or (s.data.subtitle and s.data.subtitle.strip()) or s.data.items or s.data.points)]
    
    # --- UPDATED: Call to the new sticker function ---
    request.slides = strategically_add_image_layouts(request.slides)
    request.slides = strategically_add_sticker_layouts(request.slides)
    
    slides_to_image, index_map = identify_slides_for_imaging(request.slides)
    if slides_to_image:
        try:
            payload = ImageServiceRequest(slides=slides_to_image).dict()
            async with httpx.AsyncClient(timeout=300.0) as client:
                response = await client.post(IMAGE_SERVICE_URL, json=payload)
                response.raise_for_status()
                imaged_slides = response.json().get("slides_with_images", [])
            original_indices = list(index_map.keys())
            for i, imaged_slide_data in enumerate(imaged_slides):
                request.slides[original_indices[i]].image_base64 = imaged_slide_data.get("image_base64")
        except Exception as e:
            raise HTTPException(status_code=503, detail=f"Image Service error: {e}")

    try:
        template_path = THEME_MAP.get(request.theme, DEFAULT_TEMPLATE)
        logger.info(f"[{job_id}] Using template file: {template_path}")
        if not os.path.exists(template_path):
            logger.error(f"[{job_id}] Template file not found at {template_path}. Falling back to default.")
            template_path = DEFAULT_TEMPLATE
            if not os.path.exists(template_path):
                raise HTTPException(status_code=500, detail="Default template file not found.")
        
        prs = Presentation(template_path)
        
        for slide_request in request.slides:
            layout_index = LAYOUT_MAP.get(slide_request.layout)
            if layout_index is None: continue
            
            if len(prs.slide_layouts) <= layout_index:
                logger.warning(f"[{job_id}] Layout index {layout_index} not found. Skipping slide.")
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            data = slide_request.data
            
            title_ph = get_placeholder(slide, PP_PLACEHOLDER.TITLE)
            if not title_ph:
                title_ph = get_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
            
            if title_ph:
                title_ph.text = data.title or ""
            else:
                logger.warning(f"[{job_id}] Layout index {layout_index} has no TITLE or CENTER_TITLE placeholder. Skipping title.")

            if layout_index == 0: # Title Slide
                subtitle_ph = get_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
                if subtitle_ph:
                    subtitle_ph.text = data.subtitle or ""
                
            elif layout_index == 1: # Bullet Points Slide
                body_ph = get_placeholder(slide, PP_PLACEHOLDER.BODY)
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()
                    bullets = [b.replace('**', '').strip() for b in (data.items or data.points or [])]
                    if bullets:
                        tf.text = bullets[0]
                        for item in bullets[1:]: tf.add_paragraph().text = item
                else:
                    logger.warning(f"[{job_id}] Layout 1 is missing a BODY placeholder.")

            elif layout_index in [2, 3]: # Image Slides
                body_ph = get_placeholder(slide, PP_PLACEHOLDER.BODY)
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()
                    bullets = [b.replace('**', '').strip() for b in (data.items or data.points or [])]
                    if bullets:
                        tf.text = bullets[0]
                        for item in bullets[1:]: tf.add_paragraph().text = item
                else:
                    logger.warning(f"[{job_id}] Layout {layout_index} is missing a BODY placeholder.")
                
                if slide_request.image_base64:
                    picture_ph = get_placeholder(slide, PP_PLACEHOLDER.PICTURE)
                    if picture_ph:
                        try:
                            image_data = base64.b64decode(slide_request.image_base64)
                            image_stream = io.BytesIO(image_data)
                            picture_ph.insert_picture(image_stream)
                        except Exception as e:
                            logger.warning(f"[{job_id}] Failed to insert image: {e}")
                    else:
                        logger.warning(f"[{job_id}] Layout {layout_index} is missing a PICTURE placeholder.")

            # --- ADDITIONS START HERE: Logic for new sticker layouts ---
            elif layout_index in [4, 5]: # Sticker Left & Sticker Right Slides
                body_ph = get_placeholder(slide, PP_PLACEHOLDER.BODY)
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()
                    bullets = [b.replace('**', '').strip() for b in (data.items or data.points or [])]
                    if bullets:
                        tf.text = bullets[0]
                        for item in bullets[1:]: tf.add_paragraph().text = item
                else:
                    logger.warning(f"[{job_id}] Layout {layout_index} is missing a BODY placeholder.")
                
                if slide_request.image_base64:
                    picture_ph = get_placeholder(slide, PP_PLACEHOLDER.PICTURE)
                    if picture_ph:
                        try:
                            image_data = base64.b64decode(slide_request.image_base64)
                            image_stream = io.BytesIO(image_data)
                            picture_ph.insert_picture(image_stream)
                        except Exception as e:
                            logger.warning(f"[{job_id}] Failed to insert sticker image: {e}")
                    else:
                        logger.warning(f"[{job_id}] Layout {layout_index} is missing a PICTURE placeholder for the sticker.")
            # --- ADDITIONS END HERE ---

    except Exception as e:
        logger.error(f"[{job_id}] PPTX generation failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PPTX generation failed: {e}")

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
            prs.save(tmp_pptx.name)
            temp_pptx_path = tmp_pptx.name
            
            pptx_filename = f"{job_id}.pptx"
            destination_blob = f"presentations/{pptx_filename}"
            download_url = upload_to_gcs(temp_pptx_path, destination_blob, job_id)
    finally:
        if 'temp_pptx_path' in locals() and os.path.exists(temp_pptx_path):
            os.remove(temp_pptx_path)
    
    encoded_url = urllib.parse.quote(download_url, safe='')
    preview_url = f"https://docs.google.com/gview?url={encoded_url}&embedded=true"

    logger.info(f"✅ [{job_id}] Process complete. Returning URLs.")
    return GenerationResponse(download_url=download_url, preview_url=preview_url)