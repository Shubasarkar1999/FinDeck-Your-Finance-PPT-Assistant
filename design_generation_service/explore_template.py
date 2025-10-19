# layout_inspector.py
# Utility script to visually inspect PowerPoint template layouts and placeholders.

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# --- Configuration ---
# IMPORTANT: Change this path to test a specific template file
# Example: 'templates/Dark.pptx', 'templates/Streamline.pptx', etc.
TEMPLATE_FILE = 'templates/Dark.pptx'
OUTPUT_FILE = 'layout_visual_report.pptx'

# --- Script ---
print(f" Loading template: {TEMPLATE_FILE}")
prs = Presentation(TEMPLATE_FILE)

# Create a new presentation for the visual report
report_prs = Presentation()
print(" Creating a visual report of all available layouts...\n")

for i, layout in enumerate(prs.slide_layouts):
    print(f" - Processing layout index: {i} ({len(layout.placeholders)} placeholders)")

    # Add a slide using the current layout
    slide = report_prs.slides.add_slide(layout)

    # Loop through placeholders and label them
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        ph_idx = shape.placeholder_format.idx

        try:
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                shape.text = f"Layout Index: {i}\nTitle Placeholder (ID: {ph_idx})"
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                shape.text = f"Subtitle Placeholder (ID: {ph_idx})"
            elif ph_type == PP_PLACEHOLDER.BODY:
                shape.text_frame.text = f"Body Placeholder (ID: {ph_idx})\n- For bullets and text"
            elif ph_type == PP_PLACEHOLDER.PICTURE:
                # Add a visible textbox label in place of picture placeholder
                sp = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                sp.text = f"Picture Placeholder (ID: {ph_idx})"
            else:
                if shape.has_text_frame:
                    shape.text = f"Other Placeholder (ID: {ph_idx}, Type: {ph_type})"
        except Exception as e:
            print(f"   ⚠️ Error labeling placeholder {ph_idx}: {e}")

print(f"\n Report finished! Open '{OUTPUT_FILE}' to see all your layouts.")
report_prs.save(OUTPUT_FILE)
